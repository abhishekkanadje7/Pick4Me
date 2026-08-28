import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Slate, Forest Teal, Warm Amber, Crisp White, Light Gray
    C_SLATE = RGBColor(15, 23, 42)      # #0F172A
    C_TEAL = RGBColor(13, 148, 136)     # #0D9488
    C_TEAL_DARK = RGBColor(17, 94, 89)  # #115E59
    C_AMBER = RGBColor(217, 119, 6)     # #D97706
    C_WHITE = RGBColor(255, 255, 255)
    C_LIGHT_BG = RGBColor(248, 250, 252) # #F8FAFC
    C_CARD_BG = RGBColor(241, 245, 249)  # #F1F5F9
    C_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    C_BORDER = RGBColor(226, 232, 240)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="YUVA YODHA ENERGY TECH HACKATHON 2026"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = C_TEAL
        p_c.font.name = "Arial"

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = C_SLATE
        p_t.font.name = "Arial"

        # Subtle divider
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = C_BORDER
        line.line.fill.background()

    def add_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Dark Premium Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_SLATE
    bg1.line.fill.background()

    # Track badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.2), Inches(5.8), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 41, 59)
    badge.line.color.rgb = C_TEAL
    badge.line.width = Pt(1)
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "⚡ YUVA YODHA ENERGY TECH HACKATHON 2026"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = C_TEAL
    p_b.alignment = PP_ALIGN.CENTER

    # Title & Subtitle
    t_box = s1.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Pick4Me"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Zero-Emission Hyperlocal Logistics & Smart Campus Supply Network"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_TEAL
    p2.font.name = "Arial"

    p3 = tf.add_paragraph()
    p3.text = "Empowering Local Campus Commerce & Everyday Commuters via 2-Stage Escrow Protocols"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.font.name = "Arial"

    # Metadata Cards at bottom
    meta1 = add_card(s1, Inches(0.9), Inches(4.7), Inches(3.6), Inches(1.8), RGBColor(30, 41, 59), C_TEAL)
    tf_m1 = meta1.text_frame
    tf_m1.word_wrap = True
    p = tf_m1.paragraphs[0]
    p.text = "🎯 CHALLENGE TRACK"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_AMBER
    p = tf_m1.add_paragraph()
    p.text = "Smart Buildings &\nClean Intra-Campus Mobility"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    meta2 = add_card(s1, Inches(4.8), Inches(4.7), Inches(3.6), Inches(1.8), RGBColor(30, 41, 59), C_TEAL)
    tf_m2 = meta2.text_frame
    tf_m2.word_wrap = True
    p = tf_m2.paragraphs[0]
    p.text = "🌱 CORE SUSTAINABILITY IMPACT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_TEAL
    p = tf_m2.add_paragraph()
    p.text = "100% Zero-Emission Deliveries\nVia Pedestrian Commuter Foot-Traffic"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    meta3 = add_card(s1, Inches(8.7), Inches(4.7), Inches(3.7), Inches(1.8), RGBColor(30, 41, 59), C_TEAL)
    tf_m3 = meta3.text_frame
    tf_m3.word_wrap = True
    p = tf_m3.paragraphs[0]
    p.text = "👥 PROJECT LEAD"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_AMBER
    p = tf_m3.add_paragraph()
    p.text = "Abhishek Kanadje & Team\nEmail: abhishekkanadje7@gmail.com\nLive: pick4me-l868.onrender.com"
    p.font.size = Pt(11)
    p.font.color.rgb = C_WHITE

    # ==========================================
    # SLIDE 2: THE PROBLEM (ENERGY & EMISSION CRISIS)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "The Problem: Inefficient Last-Mile Logistics in Campuses & Buildings")

    p_cards = [
        ("🚗 High Fossil Fuel Consumption", 
         "Commercial delivery services (Swiggy, Zepto, Blinkit) deploy dedicated petrol two-wheelers for tiny, low-value items (stationery, printouts, snacks, medicines).\n\nSingle-item vehicular trips create excessive fuel burning and carbon footprints per delivery.",
         C_SLATE, RGBColor(254, 242, 242)),
        ("🏭 Pollution & Traffic Congestion",
         "Motorized delivery fleets create severe noise and air pollution inside residential university hostels, academic zones, and smart building townships.\n\nMost commercial delivery agents are blocked at campus entrance gates, failing room-to-room delivery.",
         C_SLATE, RGBColor(255, 251, 235)),
        ("💸 Excessive Surge & Invisibility",
         "High minimum order limits and surge delivery fees penalize students buying small daily essentials.\n\nLocal on-campus vendors and stationery canteens remain invisible without an affordable digital storefront.",
         C_SLATE, RGBColor(240, 253, 250))
    ]

    for i, (title, desc, tc, bg) in enumerate(p_cards):
        c = add_card(s2, Inches(0.8 + i*3.95), Inches(2.0), Inches(3.8), Inches(4.8), bg, C_BORDER)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_SLATE
        p = tf.add_paragraph()
        p.text = "\n" + desc
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 3: OUR VISION & ENERGY TECH ALIGNMENT
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Our Vision: Decarbonizing Campus Mobility via Crowdsourcing")

    c_left = add_card(s3, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.8), C_CARD_BG)
    tf_l = c_left.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "🌱 The Green Logistics Hypothesis"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_TEAL
    p = tf_l.add_paragraph()
    p.text = "\nEvery day, thousands of students, researchers, and campus residents walk between gates, libraries, academic complexes, and hostels with empty bags.\n\nInstead of burning fuel with new delivery vehicles, why not monetize these existing daily pedestrian walking routes?\n\nPick4Me harnesses this natural human mobility to achieve a 100% zero-emission micro-delivery network."
    p.font.size = Pt(13)
    p.font.color.rgb = C_SLATE

    c_right = add_card(s3, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8), RGBColor(240, 253, 250), C_TEAL)
    tf_r = c_right.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "⚡ Alignment with Smart Buildings & Energy"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_TEAL_DARK
    
    points = [
        ("Zero Capital & Vehicle Emissions", "Eliminates fuel-burning transport fleets from the campus footprint."),
        ("Micro-Supply Chain for Smart Townships", "Connects local campus retail with residential hostel rooms seamlessly."),
        ("Energy Efficiency per Delivery", "Energy expenditure = 0 kWh fossil energy (pure human routine transit)."),
        ("Circular Campus Economy", "Keeps commerce and delivery rewards circulating within the student community.")
    ]
    for pt, pd in points:
        p = tf_r.add_paragraph()
        p.text = f"• {pt}: {pd}"
        p.font.size = Pt(12)
        p.font.color.rgb = C_SLATE

    # ==========================================
    # SLIDE 4: THE PICK4ME ECOSYSTEM
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "The Pick4Me Solution: 3-Party Peer Commerce & Logistics")

    roles = [
        ("🛍️ Customer A (Buyer)", "• Places orders from campus stores or posts custom errands\n• Pays 100% advance into secure Admin Escrow\n• Holds private 4-digit Delivery Handover OTP", C_TEAL),
        ("🏪 Store Merchant (Shop)", "• Catalogs stationery, lab manuals, canteens & pharmacy\n• Hands over item to Customer B at shop counter\n• Receives 100% item cost automatically into Store Wallet", C_AMBER),
        ("🚴 Customer B (Commuter)", "• Browses 'Deliver & Earn' board along their daily walking route\n• Claims delivery, collects from store, delivers to room\n• Receives delivery reward (+ Speed Bonus) into personal Wallet", C_TEAL),
        ("🛡️ Central Admin Escrow", "• Holds advance funds securely via UPI Gateway\n• Executes 2-Stage automated autonomous payouts\n• Monitors KYC verification and campus trust metrics", C_SLATE)
    ]

    for i, (rtitle, rdesc, color) in enumerate(roles):
        c = add_card(s4, Inches(0.8 + i*2.95), Inches(2.0), Inches(2.85), Inches(4.8), C_WHITE, color)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rtitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = "\n" + rdesc
        p.font.size = Pt(12)
        p.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 5: HOW IT WORKS (2-STAGE ESCROW FLOW)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "End-to-End Workflow: 2-Stage Autonomous Escrow Settlement")

    steps = [
        ("1. Order & Escrow", "Customer A orders lab manual & pays ₹120 item + ₹35 reward = ₹155 to Admin Escrow."),
        ("2. Commuter Claim", "Customer B (fellow student walking near shop) claims task on 'Deliver & Earn' board."),
        ("3. Stage 1 Handover", "Customer B picks up item at store. Shopkeeper confirms & ₹120 is credited to Shopkeeper Wallet!"),
        ("4. Stage 2 Handover", "Customer B delivers to Customer A room. Customer A shares 4-digit OTP & ₹35 (+Bonus) is credited to Customer B Wallet!")
    ]

    for i, (title, desc) in enumerate(steps):
        top_y = 2.0 + i*1.2
        c = add_card(s5, Inches(0.8), Inches(top_y), Inches(11.7), Inches(1.05), C_CARD_BG, C_TEAL)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"STEP {i+1}: {title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEAL_DARK
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_SLATE

    # ==========================================
    # SLIDE 6: CORE INNOVATIONS & SMART SLA ENGINE
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Core Technological & Operational Innovations")

    inno_cards = [
        ("📍 Dynamic Distance Tiers", 
         "Fair, automated base reward calculations across campus zones:\n• 0-1 km: ₹20 Base (25m Target SLA)\n• 1-2.5 km: ₹35 Base (35m Target SLA)\n• 2.5-5 km+: ₹55 Base (50m Target SLA)",
         C_TEAL),
        ("⚡ Smart Speed Bonus Engine", 
         "Incentivizes rapid peer handovers:\n• Superfast Delivery (< 65% SLA): +₹10 Speed Bonus credited automatically!\n• Late Delivery (> 140% SLA): -₹10 Late Penalty\n• Urgent Rush Mode (+₹15 Priority Surge)",
         C_AMBER),
        ("🛡️ 2-Stage Handover OTP", 
         "100% Anti-Fraud Protection:\n• Eliminates order cancellations & disputes\n• Funds released only upon physical item delivery verification via 4-digit private OTP",
         C_TEAL),
        ("💰 Real-Time Wallet Settlement", 
         "Autonomous transaction ledger with instant balance credit for merchants and student commuters, with zero manual intervention.",
         C_SLATE)
    ]

    for i, (ititle, idesc, color) in enumerate(inno_cards):
        x = 0.8 + (i % 2) * 5.95
        y = 2.0 + (i // 2) * 2.45
        c = add_card(s6, Inches(x), Inches(y), Inches(5.75), Inches(2.25), C_WHITE, color)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ititle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = "\n" + idesc
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_SLATE

    # ==========================================
    # SLIDE 7: TECH STACK & PRODUCTION ARCHITECTURE
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Technical Stack & Robust Production Architecture")

    techs = [
        ("Frontend Layer", "HTML5, CSS3 Variables, Flexbox/Grid, Vanilla JS, Micro-Animations & Scroll Reveal"),
        ("Backend Core", "Python 3.10+, Flask 3.0 Web Framework, Jinja2 Templating Engine, REST API architecture"),
        ("Data Persistence", "SQLite3 with Foreign Key Constraints, Transactional Integrity & Ledger Consistency"),
        ("Security & Auth", "PBKDF2:SHA256 Password Hashing, 2-Stage 4-Digit OTP, Escrow Fund Isolation"),
        ("Payment & Gateway", "UPI QR Code & VPA Gateway Integration (7387157739@upi)"),
        ("Cloud & Deployment", "Gunicorn WSGI Server, Render Cloud Service, Git CI/CD automated pipeline")
    ]

    for i, (tname, tdesc) in enumerate(techs):
        x = 0.8 + (i % 2) * 5.95
        y = 2.0 + (i // 2) * 1.6
        c = add_card(s7, Inches(x), Inches(y), Inches(5.75), Inches(1.4), C_CARD_BG, C_TEAL)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💻 {tname}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEAL_DARK
        p = tf.add_paragraph()
        p.text = tdesc
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_SLATE

    # ==========================================
    # SLIDE 8: LIVE DEPLOYED PLATFORM SHOWCASE
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Live Working Platform Showcase & Verification")

    screens = [
        ("🛍️ Amazon-Style Storefront", "Real-time search, filters, category tabs, and interactive shopping cart."),
        ("🚴 'Deliver & Earn' Board", "Live marketplace for student commuters to claim nearby deliveries on route."),
        ("🛡️ 2-Stage Escrow Timeline", "Visual milestone tracker showing Shop Pickup & Customer Delivery stages."),
        ("💰 Autonomous Wallet Ledger", "Complete financial history, payout breakdown, and instant withdrawal.")
    ]

    for i, (sname, sdesc) in enumerate(screens):
        c = add_card(s8, Inches(0.8 + i*2.95), Inches(2.0), Inches(2.85), Inches(4.8), C_WHITE, C_TEAL)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sname
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_SLATE
        p = tf.add_paragraph()
        p.text = "\n" + sdesc + "\n\n✓ 100% Automated Unit & Integration Tests Passed (`test_workflow.py`)"
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 9: ENVIRONMENTAL & ENERGY IMPACT
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Measurable Energy Savings & Environmental Impact")

    metrics = [
        ("0.15 kg", "CO2 Saved Per Delivery", "Compared to standard petrol two-wheeler deliveries (1.5 km trip).", C_TEAL),
        ("54.75 Tons", "Annual CO2 Avoided", "Per average university campus with 1,000 daily crowdsourced peer errands.", C_AMBER),
        ("18,000+ Liters", "Petrol Conserved", "Zero fossil fuel consumed for intra-campus and township deliveries.", C_TEAL),
        ("100%", "Zero Extra Vehicles", "Zero fleet congestion or noise pollution inside academic & residential premises.", C_SLATE)
    ]

    for i, (val, title, desc, color) in enumerate(metrics):
        c = add_card(s9, Inches(0.8 + i*2.95), Inches(2.0), Inches(2.85), Inches(4.8), C_CARD_BG, color)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_SLATE
        p = tf.add_paragraph()
        p.text = "\n" + desc
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 10: BUSINESS MODEL & SCALABILITY
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Business Model, Monetization & Scalability")

    b_cards = [
        ("💰 Revenue Model", 
         "• 5% Platform Convenience Fee on delivery rewards\n• Promoted store listings for local merchants\n• Zero vehicle asset overhead or maintenance costs"),
        ("📈 High Campus Scalability", 
         "• Replicable across 500+ Indian university campuses, engineering colleges, and gated corporate parks\n• Requires zero physical infrastructure or vehicle capex"),
        ("🤝 Community Multiplier", 
         "• Empowers local small businesses to compete with quick-commerce giants\n• Enables students to earn flexible micro-incomes on routine routes")
    ]

    for i, (btitle, bdesc) in enumerate(b_cards):
        c = add_card(s10, Inches(0.8 + i*3.95), Inches(2.0), Inches(3.8), Inches(4.8), C_WHITE, C_TEAL)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = btitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_SLATE
        p = tf.add_paragraph()
        p.text = "\n" + bdesc
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 11: IMPLEMENTATION ROADMAP
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Implementation & Future Development Roadmap")

    phases = [
        ("Phase 1: Foundation (COMPLETED & LIVE)", "• Full 3-Party Peer Marketplace & Store Cataloging\n• 2-Stage Automated Escrow Handover & 4-Digit OTP\n• Dynamic Distance & Time SLA Speed Bonus Engine\n• Production Cloud Deployment on Render"),
        ("Phase 2: Smart Campus IoT Integration", "• Smart QR Lockers in hostel lobbies for asynchronous drop-off\n• WhatsApp / Telegram Bot notification alerts\n• AI-assisted commuter route matching & batch errands"),
        ("Phase 3: Multi-Institution Scaling", "• Expansion to 50+ university campuses and tech hubs\n• Institutional carbon credit tracking & sustainability reporting")
    ]

    for i, (pname, pdesc) in enumerate(phases):
        top_y = 2.0 + i*1.6
        c = add_card(s11, Inches(0.8), Inches(top_y), Inches(11.7), Inches(1.4), C_CARD_BG, C_TEAL if i==0 else C_BORDER)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pname
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEAL_DARK if i==0 else C_SLATE
        p = tf.add_paragraph()
        p.text = pdesc
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_SLATE

    # ==========================================
    # SLIDE 12: TEAM & CONCLUSION
    # ==========================================
    s12 = prs.slides.add_slide(blank_layout)
    bg12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = C_SLATE
    bg12.line.fill.background()

    c_lead = add_card(s12, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8), RGBColor(30, 41, 59), C_TEAL)
    tf_l = c_lead.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Pick4Me — Transforming Everyday Commutes into Zero-Emission Deliveries"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    p = tf_l.add_paragraph()
    p.text = "\n👥 Team Members:\n• Abhishek Kanadje (Project Lead & Full-Stack Architect)\n• Email: abhishekkanadje7@gmail.com\n\n🌐 Live Application & Source Code Links:\n• Live Application Demo: https://pick4me-l868.onrender.com\n• GitHub Repository: https://github.com/abhishekkanadje7/Pick4Me\n\nThank you! We welcome your questions and feedback."
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(226, 232, 240)

    # Save presentation
    output_pptx = "/Users/abhishekkanadje7/Desktop/Pick4Me_EnergyTech_Presentation.pptx"
    prs.save(output_pptx)
    print(f"Presentation saved successfully to: {output_pptx}")

if __name__ == '__main__':
    create_deck()
